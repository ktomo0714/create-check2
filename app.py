import streamlit as st
import openai
import os
import sqlite3
import hashlib
import datetime
import io
from dotenv import load_dotenv

# 追加ライブラリ
import docx  # Word文書処理用
import pptx  # PowerPoint処理用
import PyPDF2  # PDF処理用（新規追加）

# OpenAI APIキーの設定
api_key = os.environ.get("OPENAI_API_KEY")
if not api_key:
    st.error("OpenAI APIキーが設定されていません。Renderのダッシュボードで環境変数を設定してください。")
    st.stop()

# OpenAIクライアントの初期化
import httpx
http_client = httpx.Client()
client = openai.OpenAI(api_key=api_key, http_client=http_client)

# SQLiteデータベースのセットアップ
def init_db():
    conn = sqlite3.connect('app_data.db')
    c = conn.cursor()
    
    # ユーザーテーブルの作成（存在しない場合）
    c.execute('''
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        password_hash TEXT NOT NULL,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    ''')
    
    # 履歴テーブルの作成（存在しない場合）
    c.execute('''
    CREATE TABLE IF NOT EXISTS history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        action_type TEXT NOT NULL,
        content TEXT,
        result TEXT,
        file_name TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (user_id) REFERENCES users (id)
    )
    ''')
    
    conn.commit()
    conn.close()

# パスワードハッシュ化関数
def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

# ユーザー登録関数
def register_user(username, password):
    conn = sqlite3.connect('app_data.db')
    c = conn.cursor()
    
    try:
        hashed_password = hash_password(password)
        c.execute("INSERT INTO users (username, password_hash) VALUES (?, ?)", 
                 (username, hashed_password))
        conn.commit()
        return True
    except sqlite3.IntegrityError:
        # ユーザー名が既に存在する場合
        return False
    finally:
        conn.close()

# ユーザー認証関数
def authenticate_user(username, password):
    conn = sqlite3.connect('app_data.db')
    c = conn.cursor()
    
    hashed_password = hash_password(password)
    c.execute("SELECT id FROM users WHERE username = ? AND password_hash = ?", 
             (username, hashed_password))
    user = c.fetchone()
    conn.close()
    
    if user:
        return user[0]  # ユーザーIDを返す
    else:
        return None

# 履歴を保存する関数
def save_history(user_id, action_type, content, result, file_name=None):
    conn = sqlite3.connect('app_data.db')
    c = conn.cursor()
    
    c.execute("""
    INSERT INTO history (user_id, action_type, content, result, file_name) 
    VALUES (?, ?, ?, ?, ?)
    """, (user_id, action_type, content, result, file_name))
    
    conn.commit()
    conn.close()

# ユーザーの履歴を取得する関数
def get_user_history(user_id):
    conn = sqlite3.connect('app_data.db')
    c = conn.cursor()
    
    c.execute("""
    SELECT action_type, content, result, file_name, created_at 
    FROM history 
    WHERE user_id = ? 
    ORDER BY created_at DESC
    """, (user_id,))
    
    history = c.fetchall()
    conn.close()
    
    return history

# ファイルからテキストを抽出する関数
def extract_text_from_file(uploaded_file):
    file_type = uploaded_file.name.split('.')[-1].lower()
    text = ""
    file_content = uploaded_file.getvalue()
    
    try:
        if file_type == 'txt':
            # テキストファイルの処理
            text = uploaded_file.getvalue().decode('utf-8')
        
        elif file_type in ['docx', 'doc']:
            # Wordファイルの処理
            doc = docx.Document(io.BytesIO(file_content))
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
        
        elif file_type in ['pptx', 'ppt']:
            # PowerPointファイルの処理
            prs = pptx.Presentation(io.BytesIO(file_content))
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        elif file_type == 'pdf':
            # PDFファイルの処理（新規追加）
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            
            # PDFの各ページからテキストを抽出
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n"
        
        else:
            text = "サポートされていないファイル形式です。"
    
    except Exception as e:
        text = f"ファイルの処理中にエラーが発生しました: {str(e)}"
    
    return text

# データベースの初期化
init_db()

# アプリのタイトルとスタイル
st.set_page_config(
    page_title="生成・校閲アプリケーション",
    page_icon="📝",
    layout="wide"
)

# セッション状態の初期化
if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
if 'user_id' not in st.session_state:
    st.session_state.user_id = None
if 'username' not in st.session_state:
    st.session_state.username = None

# ログイン機能
def login_page():
    st.title("ログイン")
    
    tab1, tab2 = st.tabs(["ログイン", "新規登録"])
    
    with tab1:
        username = st.text_input("ユーザー名", key="login_username")
        password = st.text_input("パスワード", type="password", key="login_password")
        
        if st.button("ログイン", key="login_button"):
            if username and password:
                user_id = authenticate_user(username, password)
                if user_id:
                    st.session_state.logged_in = True
                    st.session_state.user_id = user_id
                    st.session_state.username = username
                    st.success("ログインに成功しました！")
                    st.rerun()
                else:
                    st.error("ユーザー名またはパスワードが間違っています。")
            else:
                st.warning("ユーザー名とパスワードを入力してください。")
    
    with tab2:
        new_username = st.text_input("新しいユーザー名", key="register_username")
        new_password = st.text_input("新しいパスワード", type="password", key="register_password")
        confirm_password = st.text_input("パスワード（確認）", type="password", key="confirm_password")
        
        if st.button("登録", key="register_button"):
            if new_username and new_password and confirm_password:
                if new_password == confirm_password:
                    if register_user(new_username, new_password):
                        st.success("アカウントが作成されました。ログインしてください。")
                    else:
                        st.error("そのユーザー名は既に使用されています。")
                else:
                    st.error("パスワードが一致しません。")
            else:
                st.warning("すべての項目を入力してください。")

# サイドバーメニュー
def sidebar_menu():
    with st.sidebar:
        st.title(f"こんにちは、{st.session_state.username}さん")
        
        st.title("機能選択")
        app_mode = st.radio(
            "モードを選択してください:",
            ["テキスト生成", "テキスト校閲", "履歴閲覧"],
            label_visibility="visible"  # ラベルを表示する
        )
        
        st.divider()
        
        # APIモデル選択
        model = st.selectbox(
            "使用するモデル:",
            ["gpt-4o-mini","gpt-4o"],
            index=0
        )
        
        # 温度設定（クリエイティビティの調整）
        temperature = st.slider(
            "温度 (クリエイティビティ)", 
            0.0, 1.0, 0.3, 0.1,
            label_visibility="visible"  # ラベルを表示する
        )
        
        st.divider()
        st.write("生成・校閲アプリケーション")
        
        if st.button("ログアウト"):
            st.session_state.logged_in = False
            st.session_state.user_id = None
            st.session_state.username = None
            st.rerun()
            
        return app_mode, model, temperature

# メイン関数
def main():
    # ログイン状態の確認
    if not st.session_state.logged_in:
        login_page()
    else:
        # ログイン済みの場合、メイン機能を表示
        app_mode, model, temperature = sidebar_menu()
        
        st.title("生成・校閲アプリケーション")
        
        if app_mode == "テキスト生成":
            text_generation(model, temperature)
        elif app_mode == "テキスト校閲":
            text_proofreading(model, temperature)
        elif app_mode == "履歴閲覧":
            view_history()

# テキスト生成機能
def text_generation(model, temperature):
    st.header("テキスト生成")
    
    prompt_type = st.selectbox(
        "生成するテキストのタイプ:",
        ["メールマガジン", "SMS", "SNS投稿"]
    )
    
    topic = st.text_input("トピックや主題:")
    
    length = st.select_slider(
        "文章の長さ:",
        options=["短め (100字程度)", "標準 (300字程度)", "長め (500字程度)", "詳細 (1000字以上)"]
    )
    
    additional_info = st.text_area("追加情報や要望があれば入力してください:")
    
    if st.button("生成する", type="primary"):
        if not topic:
            st.warning("トピックを入力してください。")
        else:
            with st.spinner("AIが文章を生成中..."):
                prompt = f"""
                次の条件に合うテキストを生成してください:
                - タイプ: {prompt_type}
                - トピック: {topic}
                - 長さ: {length}
                - 追加情報: {additional_info}
                
                日本語で自然な文章を生成してください。
                あなたは金融機関の広告作成を専門的に支援するAIアシスタントです。金融商品・サービスの広告作成において、法令遵守と効果的なコミュニケーションを両立させる提案を行います。

                ## 基本方針
                - 金融商品取引法、銀行法、保険業法など関連法規に完全準拠した広告コンテンツを生成する
                - 誤解を招く表現や過度な期待を抱かせる表現を徹底的に排除する
                - リスクとリターンの適切なバランスを保った説明を心がける
                - 対象顧客層に応じた適切な表現と情報量を選択する
                - 金融機関としての信頼性・安定性を表現しつつ、差別化ポイントを明確に伝える

                ## 広告種類別のガイドライン
                    **Web広告・バナー**
                     - 簡潔で明確なメッセージと視覚的一貫性
                     - クリック後のランディングページとの整合性
                     - 小さなスペースでも必要な免責事項を表示
                     - CTAの明確さと行動喚起の適切さ

                    **パンフレット・商品説明資料**
                    - 段階的な情報提供による理解促進
                    - 重要事項の視認性確保
                    - 図表・イラストの効果的活用
                    - 商品構造・手数料体系の透明な説明

                    **ソーシャルメディア投稿**
                    - プラットフォーム特性に合わせた最適な表現
                    - エンゲージメントと法令遵守のバランス
                    - シリーズ投稿による段階的な情報提供
                    - コメント対応のための想定Q&A

                ## コンプライアンス要件
                    **必須開示事項**
                    - 金融機関名・登録番号
                    - 手数料・費用の明示
                    - リスク情報の適切な開示
                    - 実績数値使用時の出典・条件明示

                    **禁止表現**
                    - 元本保証がない商品の「安全」「確実」等の表現
                    - 利回り・リターンの断定的表現
                    - 他社比較における不適切な優位性主張
                    - 顧客の投資判断を誤らせる表現

                    **適正表示**
                    - リスク文言の視認性（文字サイズ、表示時間等）
                    - 条件付き表現の条件明示
                    - 専門用語の平易な説明
                    - 図表・グラフの適切な縮尺と説明

                ## 広告効果向上のポイント
                    **ターゲティング**
                    - 顧客セグメント別のニーズ・関心事への合致
                    - 金融リテラシーレベルに応じた表現の選択
                    - ライフイベントに合わせたメッセージング
                    - 商品特性と顧客属性のマッチング

                    **差別化要素**
                    - 金利・手数料等の定量的優位性
                    - サービス・サポートの質的優位性
                    - テクノロジー・利便性の革新性
                    - 社会的意義・ESG要素の訴求

                    **心理的アプローチ**
                    - 安心感・信頼性の醸成
                    - 将来不安の解消・目標達成の支援
                    - 社会的証明による後押し
                    - 希少性・適時性の適切な強調

                ## 生成プロセス
                    1. 広告目的と対象商品・サービスの明確化
                    2. ターゲット顧客層と媒体の特定
                    3. 主要メッセージと差別化ポイントの設定
                    4. コンプライアンス要件の確認とリスク開示の組み込み
                    5. 広告クリエイティブの生成（複数バージョン）
                    6. コンプライアンス最終チェック

                ## 注意事項
                    - 投資・保険商品の広告はとりわけ厳格な規制があることを常に意識する
                    - 広告表現の解釈は多様であることを考慮し、慎重な表現選択を行う
                    - ハルシネーション（誤った情報の生成）を防止し、不確かな内容は含めない
                    - 最新の金融規制に基づいた広告表現を心がけ、必要に応じて確認を促す
                    - 生成した広告案は必ず金融機関のコンプライアンス部門の確認を受けるよう注記する
                    効果的な訴求と厳格なコンプライアンス準拠を両立し、金融機関と顧客双方の価値を高める広告制作を支援します。
                """
                
                try:
                    response = client.chat.completions.create(
                        model=model,
                        messages=[{"role": "user", "content": prompt}],
                        temperature=temperature,
                    )
                    
                    result = response.choices[0].message.content
                    
                    # 履歴に保存
                    save_history(
                        st.session_state.user_id, 
                        "テキスト生成", 
                        prompt, 
                        result
                    )
                    
                    st.success("テキストが生成されました！")
                    st.text_area(
                        label="生成されたテキスト:", 
                        value=result, 
                        height=300, 
                        key="generated_text",
                        label_visibility="visible"  # ラベルを表示する
                    )
                    
                    # ダウンロードボタン
                    st.download_button(
                        label="テキストをダウンロード",
                        data=result,
                        file_name=f"{topic}_generated_text.txt",
                        mime="text/plain"
                    )
                
                except Exception as e:
                    st.error(f"エラーが発生しました: {str(e)}")

# テキスト校閲機能
def text_proofreading(model, temperature):
    st.header("テキスト校閲")
    
    # 入力方式の選択
    input_method = st.radio(
        "入力方式を選択してください:", 
        ["テキスト直接入力", "ファイルアップロード"],
        label_visibility="visible"  # ラベルを表示する
    )
    
    input_text = ""
    file_name = None
    
    if input_method == "テキスト直接入力":
        input_text = st.text_area(
            label="校閲したいテキストを入力してください:", 
            height=200, 
            key="proofread_input",
            label_visibility="visible"  # ラベルを表示する
        )
    else:
        st.write("対応ファイル形式: .txt, .doc, .docx, .ppt, .pptx, .pdf")  # PDFを追加
        uploaded_file = st.file_uploader(
            label="ファイルをアップロードしてください", 
            type=["txt", "doc", "docx", "ppt", "pptx", "pdf"],
            label_visibility="visible"  # ラベルを表示する
        )  # PDFを追加
        
        if uploaded_file is not None:
            file_name = uploaded_file.name
            
            # ファイル情報の表示
            file_details = {
                "ファイル名": uploaded_file.name,
                "ファイルタイプ": uploaded_file.type,
                "ファイルサイズ": f"{uploaded_file.size} bytes"
            }
            st.write("**ファイル情報:**")
            for k, v in file_details.items():
                st.write(f"- {k}: {v}")
            
            # ファイルからテキストを抽出
            with st.spinner("ファイルからテキストを抽出中..."):
                input_text = extract_text_from_file(uploaded_file)
                
                if input_text:
                    st.subheader("抽出されたテキスト:")
                    st.write(input_text[:1000] + ("..." if len(input_text) > 1000 else ""))
                    
                    # テキスト全体の表示トグル
                    if len(input_text) > 1000:
                        with st.expander("テキスト全体を表示"):
                            st.text_area(
                                label="テキスト全体", 
                                value=input_text, 
                                height=300, 
                                key="full_extracted_text",
                                label_visibility="collapsed"  # ラベルを非表示（存在するが表示しない）
                            )
                else:
                    st.error("テキストを抽出できませんでした。")
    
    check_options = st.multiselect(
        "確認項目:",
        ["景品表示法への抵触がないか","金融商品取引法への抵触がないか","文法/スペル","わかりやすさ", "一貫性"]
    )
    
    if st.button("校閲する", type="primary"):
        if not input_text:
            st.warning("テキストを入力またはファイルをアップロードしてください。")
        else:
            with st.spinner("AIが校閲中..."):
                checks = ", ".join(check_options) if check_options else "すべての側面"
                
                prompt = f"""
                あなたは金融機関の文書校閲を専門とするAIアシスタントです。正確で信頼性の高い校閲サービスを提供し、金融業界特有の表現、規制要件、コンプライアンスを考慮した適切な修正提案を行います。
                以下のテキストを校閲してください。{checks}に注目して改善点を指摘し、
                修正案を提案してください。元のテキストを尊重しつつ、より明確で効果的な表現を目指してください。

                ## 校閲の基本方針
                    - 金融関連法規制に準拠した表現であるかを厳格に確認する
                    - 数値、金額、日付、商品名等の正確性を最優先で確認する
                    - 専門用語と平易な表現のバランスを適切に保つ
                    - 表現の一貫性と統一性を確保する
                    - リスク開示が適切かつ十分であるかを確認する
                    - わかりやすさと正確さを両立した文章構成を心がける

                ## 校閲対象文書
                    **顧客向け資料**
                    - 金融商品説明資料・パンフレット
                    - 契約書・約款
                    - 重要事項説明書
                    - 顧客宛て通知文

                    **内部文書**
                    - 業務マニュアル・手順書
                    - 社内報告書・提案書
                    - 社内規程・ポリシー
                    - 研修資料

                    **公開文書**
                    - プレスリリース
                    - IR資料・ディスクロージャー
                    - 採用情報・企業案内
                    - ウェブサイトコンテンツ

                ## 校閲のポイント
                    **法令遵守の観点**
                    - 誤解を招く表現や断定的な表現の排除
                    - 優位性を示す表現の適切性確認
                    - 必要な免責事項・注意書きの確認
                    - 個人情報保護に関する表現の確認

                    **表現・用語の観点**
                    - 専門用語の適切な使用と説明
                    - 敬語・謙譲語・丁寧語の正しい使用
                    - カタカナ語・外来語の統一表記
                    - 曖昧表現・冗長表現の修正

                    **構成・可読性の観点**
                    - 論理展開の一貫性と明確さ
                    - 段落構成・見出しの適切性
                    - 箇条書き・図表の効果的な活用
                    - フォントサイズ・書式の統一性

                    **金融特有の観点**
                    - リスク・リターンのバランスある説明
                    - 手数料・費用の明確な表示
                    - 数値・計算例の正確性
                    - 市場予測に関する適切な表現

                テキスト:
                {input_text}
                
                以下の形式で回答してください：
                1. 全体的な評価
                2. 具体的な改善点（元の文と修正案を対比）
                3. 修正後の全文

                ## 注意事項
                    - 内容の事実確認は行わず、表現・構成のみを校閲する
                    - 業界固有の専門用語や略語の使用については慎重に判断する
                    - ハルシネーション（誤った情報の生成）を防止し、不確かな修正は提案しない
                    - 文書の目的や対象読者を考慮した校閲を心がける
                    - 金融商品の内容自体に関する評価・判断は行わない
                """
                
                try:
                    response = client.chat.completions.create(
                        model=model,
                        messages=[{"role": "user", "content": prompt}],
                        temperature=temperature,
                    )
                    
                    result = response.choices[0].message.content
                    
                    # 履歴に保存
                    save_history(
                        st.session_state.user_id, 
                        "テキスト校閲", 
                        input_text, 
                        result,
                        file_name
                    )
                    
                    st.success("校閲が完了しました！")
                    
                    # タブで表示
                    tab1, tab2 = st.tabs(["校閲結果", "比較"])
                    with tab1:
                        st.markdown(result)
                        
                        # ダウンロードボタン
                        st.download_button(
                            label="校閲結果をダウンロード",
                            data=result,
                            file_name="proofreading_result.txt",
                            mime="text/plain"
                        )
                    
                    with tab2:
                        col1, col2 = st.columns(2)
                        with col1:
                            st.subheader("元のテキスト")
                            st.text_area(
                                label="元のテキスト", 
                                value=input_text, 
                                height=300, 
                                key="original_text_area",
                                label_visibility="collapsed"  # ラベルを非表示（存在するが表示しない）
                            )
                        with col2:
                            st.subheader("校閲後の提案")
                            # ここは実際には校閲後のテキストだけを抽出する必要があります
                            # 簡易的な実装として全体を表示
                            st.text_area(
                                label="校閲後の提案", 
                                value=result, 
                                height=300, 
                                key="proofread_text_area",
                                label_visibility="collapsed"  # ラベルを非表示（存在するが表示しない）
                            )
                
                except Exception as e:
                    st.error(f"エラーが発生しました: {str(e)}")

# 履歴閲覧機能
def view_history():
    st.header("利用履歴")
    
    history = get_user_history(st.session_state.user_id)
    
    if not history:
        st.info("まだ履歴がありません。")
    else:
        # フィルタリングオプション
        col1, col2 = st.columns([3, 1])
        with col1:
            action_filter = st.selectbox(
                "表示する操作タイプ:", 
                ["すべて", "テキスト生成", "テキスト校閲"],
                index=0
            )
        
        # ダウンロードボタンの配置
        with col2:
            # 履歴をテキスト形式に変換する関数
            def convert_history_to_text(history_data):
                text = "# 履歴一覧\n\n"
                for idx, (action_type, content, result, file_name, timestamp) in enumerate(history_data, 1):
                    text += f"## {idx}. {action_type} - {timestamp}\n"
                    if file_name:
                        text += f"ファイル名: {file_name}\n"
                    text += "\n### 入力内容\n"
                    text += f"{content}\n\n"
                    text += "### 結果\n"
                    text += f"{result}\n\n"
                    text += "---\n\n"
                return text
            
            # 履歴をCSV形式に変換する関数
            def convert_history_to_csv(history_data):
                csv_content = "No,操作タイプ,ファイル名,タイムスタンプ,入力内容,結果\n"
                for idx, (action_type, content, result, file_name, timestamp) in enumerate(history_data, 1):
                    # CSVでの特殊文字のエスケープ処理
                    safe_content = content.replace('"', '""') if content else ""
                    safe_result = result.replace('"', '""') if result else ""
                    safe_file_name = file_name.replace('"', '""') if file_name else ""
                    
                    csv_content += f'{idx},"{action_type}","{safe_file_name}","{timestamp}","{safe_content}","{safe_result}"\n'
                return csv_content
        
        # 履歴のフィルタリングと表示
        filtered_history = history
        if action_filter != "すべて":
            filtered_history = [h for h in history if h[0] == action_filter]
        
        if not filtered_history:
            st.info(f"{action_filter}の履歴はありません。")
        else:
            for i, (action_type, content, result, file_name, timestamp) in enumerate(filtered_history):
                history_title = f"{action_type} - {timestamp}"
                if file_name:
                    history_title += f" ({file_name})"
                    
                with st.expander(history_title):
                    st.subheader("入力内容")
                    st.text_area(
                        label="入力内容", 
                        value=content, 
                        height=100, 
                        key=f"content_{i}",
                        label_visibility="collapsed"  # ラベルを非表示（存在するが表示しない）
                    )
                    
                    st.subheader("結果")
                    st.text_area(
                        label="結果", 
                        value=result, 
                        height=200, 
                        key=f"result_{i}",
                        label_visibility="collapsed"  # ラベルを非表示（存在するが表示しない）
                    )
                    
                    # 個別履歴のダウンロードボタン
                    single_txt_data = f"# {action_type} - {timestamp}\n"
                    if file_name:
                        single_txt_data += f"ファイル名: {file_name}\n"
                    single_txt_data += "\n## 入力内容\n"
                    single_txt_data += f"{content}\n\n"
                    single_txt_data += "## 結果\n"
                    single_txt_data += f"{result}\n"
                    
                    st.download_button(
                        label="この履歴をダウンロード",
                        data=single_txt_data,
                        file_name=f"{action_type}_{timestamp.replace(':', '-').replace(' ', '_')}.txt",
                        mime="text/plain",
                        key=f"download_single_{i}"
                    )

# フッター
def footer():
    st.markdown("---")
    st.markdown("このアプリケーションは主としてOpenAI GPT-4o-mini APIを使用しています。生成されたテキストは参考用途にのみご利用ください。")

# アプリケーションの実行
if __name__ == "__main__":
    main()
    footer()
